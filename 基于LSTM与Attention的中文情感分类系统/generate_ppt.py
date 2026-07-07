"""
生成组员A汇报PPT — 学术论文答辩风格（中文）
"""
import os, io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

BASE = os.path.dirname(os.path.abspath(__file__))
RESULTS = os.path.join(BASE, "results")

# ═══════════════════════════════════════════════
# 配色
# ═══════════════════════════════════════════════
NAVY   = RGBColor(0x0F, 0x17, 0x2A)
BLUE   = RGBColor(0x25, 0x62, 0xEB)
LIGHT  = RGBColor(0xF1, 0xF5, 0xF9)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1E, 0x29, 0x3B)
GRAY   = RGBColor(0x64, 0x74, 0x8B)
MUTED  = RGBColor(0x94, 0xA3, 0xB8)
RED    = RGBColor(0xDC, 0x26, 0x26)
GREEN  = RGBColor(0x16, 0xA3, 0x4A)
ORANGE = RGBColor(0xEA, 0x58, 0x0C)
BG_RED   = RGBColor(0xFE, 0xF2, 0xF2)
BG_GREEN = RGBColor(0xF0, 0xFD, 0xF4)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ═══════════════════════════════════════════════
# 公式渲染
# ═══════════════════════════════════════════════
FORMULAS = {}
def _render(fid, latex, fs=15):
    fig, ax = plt.subplots(figsize=(0.01, 0.01))
    t = ax.text(0, 0, f"${latex}$", fontsize=fs, ha='left', va='bottom', math_fontfamily='stix')
    fig.canvas.draw()
    bb = t.get_window_extent(renderer=fig.canvas.get_renderer())
    w, h = bb.width / fig.dpi + 0.10, bb.height / fig.dpi + 0.08
    plt.close(fig)
    fig, ax = plt.subplots(figsize=(max(w, 1.5), max(h, 0.28)))
    ax.axis('off')
    ax.text(0, 0.02, f"${latex}$", fontsize=fs, ha='left', va='bottom', color='#1E293B', math_fontfamily='stix')
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=200, bbox_inches='tight', pad_inches=0.03, transparent=True)
    plt.close(fig); buf.seek(0)
    FORMULAS[fid] = (buf, w, h)
    return buf, w, h

def formula(slide, left, top, fid, latex, max_w=5.5, fs=15):
    if fid not in FORMULAS: _render(fid, latex, fs)
    buf, w, h = FORMULAS[fid]
    sld = slide.slide if hasattr(slide, 'slide') else slide
    pic = sld.shapes.add_picture(buf, Inches(left), Inches(top))
    sc = min(max_w / w, 1.0)
    pic.width = int(w * 914400 * sc)
    pic.height = int(pic.height * sc)
    return pic, h * sc

# ═══════════════════════════════════════════════
# 布局组件
# ═══════════════════════════════════════════════
class Slide:
    def __init__(self, num=None):
        self.slide = prs.slides.add_slide(prs.slide_layouts[6])
        self.num = num
        bg = self.slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
        bar = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1.05))
        bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
        ln = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.05), prs.slide_width, Inches(0.025))
        ln.fill.solid(); ln.fill.fore_color.rgb = BLUE; ln.line.fill.background()
        if num: T(self, 12.3, 7.05, 0.8, 0.35, str(num), 11, GRAY, PP_ALIGN.RIGHT)

    def title(self, text):
        T(self, 0.7, 0.18, 12, 0.7, text, 28, WHITE, PP_ALIGN.LEFT, bold=True)

def _set_font(para_or_run, size=None, bold=None, color=None, alignment=None):
    """设置段落/run 的 Latin + East Asian 双字体"""
    if hasattr(para_or_run, 'font'):
        rPr = para_or_run.font._rPr
    else:
        rPr = para_or_run._pPr
    # Latin font
    latin = rPr.find(qn('a:latin'))
    if latin is None:
        latin = rPr.makeelement(qn('a:latin'), {})
        rPr.insert(0, latin)
    latin.set('typeface', 'Times New Roman')
    # East Asian font
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.insert(1, ea)
    ea.set('typeface', 'SimSun')
    # 其他属性
    if size is not None:
        para_or_run.font.size = Pt(size)
    if bold is not None:
        para_or_run.font.bold = bold
    if color is not None:
        para_or_run.font.color.rgb = color

def T(slide, l, t, w, h, text, size=15, color=DARK, align=PP_ALIGN.LEFT, bold=False):
    tb = slide.slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.alignment = align
    _set_font(p, size, bold, color)
    return tf

def B(slide, l, t, w, h, text, color=BLUE, size=18):
    T(slide, l, t, w, h, text, size, color, bold=True)

def C(slide, x, y, ww, hh, border):
    s = slide.slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(ww), Inches(hh))
    s.fill.solid(); s.fill.fore_color.rgb = LIGHT
    s.line.color.rgb = border; s.line.width = Pt(1.5)
    return s

def L(slide, x, y, w):
    s = slide.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.02))
    s.fill.solid(); s.fill.fore_color.rgb = BLUE; s.line.fill.background()

def img(slide, path, l, t, w, h=None):
    p = os.path.join(BASE, path)
    if os.path.exists(p):
        if h: slide.slide.shapes.add_picture(p, Inches(l), Inches(t), Inches(w), Inches(h))
        else: slide.slide.shapes.add_picture(p, Inches(l), Inches(t), Inches(w))


# ═══════════════════════════════════════════════
# 第 1 页 · 封面
# ═══════════════════════════════════════════════
s = Slide()
bg = s.slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
T(s, 1.2, 1.8, 11, 1.5, "基于 LSTM 与 Attention 的中文情感分类系统", 38, WHITE, PP_ALIGN.CENTER, True)
T(s, 1.2, 3.5, 11, 0.5, "Chinese Sentiment Classification Based on LSTM and Attention Mechanism", 16, MUTED, PP_ALIGN.CENTER)
L(s, 5.2, 4.3, 3.0)
T(s, 1.2, 4.6, 11, 0.4, "自然语言处理  课程项目汇报", 15, RGBColor(0xCB,0xD5,0xE1), PP_ALIGN.CENTER)
T(s, 1.2, 5.2, 11, 0.4, "组员 A：数据管线构建  ·  RNN 基线模型  ·  项目综述", 15, MUTED, PP_ALIGN.CENTER)
T(s, 1.2, 6.2, 11, 0.3, "2026 年 7 月", 12, GRAY, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# 第 2 页 · 问题定义
# ═══════════════════════════════════════════════
s = Slide(1)
s.title("一、问题定义与研究目标")

B(s, 0.8, 1.35, 5, 0.35, "任务描述")
T(s, 0.8, 1.75, 11, 0.3, "给定一段中文评论文本，判断其情感极性：正面 / 负面", 15, DARK)
T(s, 0.8, 2.1, 11, 0.3, "形式化：二分类问题，y ∈ {0, 1}，0 = 负面，1 = 正面", 15, DARK)

B(s, 0.8, 2.65, 5, 0.35, "研究目标")
T(s, 0.8, 3.05, 11, 0.3, "构建完整的中文情感分析管线，对比五种模型架构的性能差异", 15, DARK)
T(s, 0.8, 3.4, 11, 0.3, "探究数据规模与模型复杂度对分类准确率的相对影响", 15, DARK)

B(s, 0.8, 3.95, 5, 0.35, "对比模型")
T(s, 0.8, 4.35, 11, 0.3, "RNN  →  BiLSTM  →  Attention-LSTM  →  CNN-BiLSTM  →  BERT（预训练基线）", 15, BLUE, bold=True)

B(s, 0.8, 4.9, 5, 0.35, "损失函数")
_, h = formula(s, 0.8, 5.4, "loss", r"\mathcal{L} = -\frac{1}{N}\sum_{i=1}^{N}\left[y_i\log\hat{p}_i + (1-y_i)\log(1-\hat{p}_i)\right]", 11.5, 14)

B(s, 0.8, 6.2, 5, 0.3, "数据规模")
T(s, 0.8, 6.55, 11, 0.3, "16K（初版）→ 80K（终版），正负各 40,000 条，覆盖酒店、外卖、电商三大领域", 14, DARK)

# ═══════════════════════════════════════════════
# 第 3 页 · 数据来源
# ═══════════════════════════════════════════════
s = Slide(2)
s.title("二、数据获取与预处理")

datasets = [
    ("ChnSentiCorp", "9,600 条", "酒店评论", "HuggingFace 镜像"),
    ("waimai_10k", "11,987 条", "外卖评论", "阿里云 OSS"),
    ("online_shopping", "62,774 条", "电商 10 类别", "本地 CSV 文件"),
]
for i, (name, n, dom, src) in enumerate(datasets):
    x = 0.8 + i*4.2
    C(s, x, 1.3, 3.8, 1.5, [BLUE, ORANGE, GREEN][i])
    T(s, x+0.25, 1.4, 3.2, 0.35, name, 17, [BLUE, ORANGE, GREEN][i], bold=True)
    T(s, x+0.25, 1.8, 3.2, 0.28, n, 14, DARK)
    T(s, x+0.25, 2.12, 3.2, 0.25, dom, 12, GRAY)
    T(s, x+0.25, 2.42, 3.2, 0.25, f"来源：{src}", 11, MUTED)

T(s, 1.5, 3.2, 10.3, 0.5, "合并 84,207 条  →  平衡采样  →  80,000 条（每类各 40,000）", 20, BLUE, PP_ALIGN.CENTER, True)

img(s, "results/analysis/sentiment_distribution.png", 3.6, 3.9, 6.0, 2.6)

B(s, 0.8, 6.7, 5, 0.3, "预处理流程")
T(s, 0.8, 7.0, 11.5, 0.25, "清洗 → jieba 分词 → 停用词过滤（保留否定词） → 词表构建（29,343 词） → Word2Vec（Skip-gram, 300d） → 前向 Padding（67 词）", 13, DARK)

# ═══════════════════════════════════════════════
# 第 4 页 · 关键改进
# ═══════════════════════════════════════════════
s = Slide(3)
s.title("三、预处理关键改进")

steps = ["原始文本", "清洗", "分词\n(jieba)", "停用词\n过滤", "词表\n(29,343)", "Word2Vec\n(300d)", "Padding\n(67)"]
for i, st in enumerate(steps):
    x = 0.6 + i*1.78
    sh = s.slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.35), Inches(1.5), Inches(0.8))
    sh.fill.solid(); sh.fill.fore_color.rgb = BLUE if i < 6 else NAVY; sh.line.fill.background()
    tf = sh.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = st; p.alignment = PP_ALIGN.CENTER
    _set_font(p, 12, True, WHITE)

# 对比框
sh = s.slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(2.55), Inches(5.8), Inches(2.5))
sh.fill.solid(); sh.fill.fore_color.rgb = BG_RED; sh.line.color.rgb = RED; sh.line.width = Pt(1.5)
T(s, 0.9, 2.65, 5.2, 0.35, "原方案：启发式单字过滤", 16, RED, bold=True)
T(s, 0.9, 3.1, 5.2, 0.28, "规则：len(w) < 2 → 直接丢弃", 14, DARK)
T(s, 0.9, 3.45, 5.2, 0.28, "关键情感词被误删：好、坏、差、不、低、贵", 14, RED)
T(s, 0.9, 3.8, 5.2, 0.28, "典型错误：", 14, DARK, bold=True)
T(s, 0.9, 4.1, 5.2, 0.3, '\u201c服务不热情\u201d \u2192 \u201c服务热情\u201d（语义反转！）', 14, RED, bold=True)
T(s, 0.9, 4.5, 5.2, 0.28, '\u201c性价比很低\u201d \u2192 \u201c性价比\u201d（否定信息丢失）', 14, RED)

sh = s.slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.55), Inches(5.9), Inches(2.5))
sh.fill.solid(); sh.fill.fore_color.rgb = BG_GREEN; sh.line.color.rgb = GREEN; sh.line.width = Pt(1.5)
T(s, 7.1, 2.65, 5.2, 0.35, "改进方案：精确停用词表", 16, GREEN, bold=True)
_, h = formula(s, 7.1, 3.2, "sw", r"W' = \{\,w \in W \mid w \notin S\,\}", 4.5, 14)
T(s, 7.1, 3.8, 5.2, 0.28, "删除：虚词（的/了/在/是）、语气词、标点、数字", 14, DARK)
T(s, 7.1, 4.15, 5.2, 0.28, "保留：否定词（不/没/非）+ 程度词（很/太/最）+ 单字实词", 14, GREEN)

T(s, 0.8, 5.35, 11.5, 0.35, "效果：所有模型准确率提升 1~4 个百分点，否定结构得到完整保留", 15, DARK, bold=True)

B(s, 0.8, 5.85, 5, 0.35, "管线参数")
T(s, 0.8, 6.25, 5.7, 0.28, "词表大小：29,343（min_freq = 3）", 14, DARK)
T(s, 0.8, 6.55, 5.7, 0.28, "词向量：Skip-gram，300 维，窗口 = 5", 14, DARK)
T(s, 0.8, 6.85, 5.7, 0.28, "序列长度：67（95% 分位数）", 14, DARK)

T(s, 6.8, 6.25, 5.7, 0.28, "训练轮数：30 epochs", 14, DARK)
T(s, 6.8, 6.55, 5.7, 0.28, "覆盖：100%（29,343 / 29,343）", 14, DARK)
T(s, 6.8, 6.85, 5.7, 0.28, "Padding：前向补零，保留末尾", 14, DARK)

# ═══════════════════════════════════════════════
# 第 5 页 · Word2Vec
# ═══════════════════════════════════════════════
s = Slide(4)
s.title("四、词向量训练：Skip-gram 与负采样")

B(s, 0.8, 1.35, 5.5, 0.35, "训练配置")
params = [("架构", "Skip-gram（sg = 1）"), ("维度", "300"),
          ("上下文窗口", "5"), ("训练轮数", "30"),
          ("最低词频", "3"), ("负采样数", "K = 5"),
          ("覆盖率", "100%（29,343 / 29,343）")]
for i, (k, v) in enumerate(params):
    T(s, 0.8, 1.8 + i*0.3, 5.5, 0.28, f"  {k}：{v}", 14, DARK)

B(s, 0.8, 4.2, 5.5, 0.35, "语义空间性质")
T(s, 0.8, 4.6, 5.5, 0.28, "语义相近的词在向量空间中距离更近", 14, DARK)
T(s, 0.8, 4.92, 5.5, 0.28, "余弦相似度可用于度量词语义相关性", 14, DARK)

# 右侧公式
B(s, 7.0, 1.35, 5.5, 0.35, "目标函数")
wy = 1.85
T(s, 7.0, wy, 5.5, 0.28, "给定中心词 c，预测上下文词 o：", 13, DARK)
wy += 0.32
_, h = formula(s, 7.0, wy, "w2v", r"\max\ \sum_{t=1}^{T}\sum_{-c\leq j\leq c,\,j\neq 0}\log P(w_{t+j}\mid w_t)", 5.5, 13)
wy += h + 0.18
T(s, 7.0, wy, 5.5, 0.28, "负采样近似（避免全词表 softmax）：", 13, DARK)
wy += 0.3
_, h = formula(s, 7.0, wy, "ns", r"\max\ \log\sigma(u_o^{\top}v_c) + \sum_{k=1}^{K}\mathbb{E}_{w_k\sim P_n}\left[\log\sigma(-u_k^{\top}v_c)\right]", 5.5, 13)
wy += h + 0.25

B(s, 7.0, wy, 5.5, 0.35, "嵌入矩阵")
T(s, 7.0, wy+0.42, 5.5, 0.28, "E ∈ R^(V × 300)，V = 29,343", 13, DARK)
T(s, 7.0, wy+0.72, 5.5, 0.28, "预训练权重作为模型 Embedding 层初始化", 13, DARK)

img(s, "results/analysis/wordclouds.png", 0.8, 5.5, 11.7, 1.7)

# ═══════════════════════════════════════════════
# 第 6 页 · RNN 模型
# ═══════════════════════════════════════════════
s = Slide(5)
s.title("五、RNN 基线模型：双向循环神经网络")

B(s, 0.8, 1.35, 5.5, 0.35, "模型架构")
arch = [
    "输入：词索引序列  (batch, 67)",
    "Embedding：预训练 300d  →  (batch, 67, 300)",
    "2 层 BiRNN：隐藏层 256 维 / 方向",
    "末层隐状态拼接  →  (batch, 512)",
    "Dropout (p = 0.5)  →  Linear (512 → 2)  →  Softmax",
]
for i, a in enumerate(arch):
    T(s, 0.8, 1.8 + i*0.27, 5.5, 0.25, a, 13, DARK)
T(s, 0.8, 3.3, 5.5, 0.3, "参数量：9,483,862（≈ 950 万）", 15, BLUE, bold=True)

B(s, 0.8, 3.85, 5.5, 0.35, "训练策略")
cfg = [
    "优化器：Adam（lr = 10⁻³, β₁ = 0.9, β₂ = 0.999）",
    "调度器：ReduceLROnPlateau（patience = 3, factor = 0.5）",
    "损失函数：交叉熵",
    "正则化：Dropout(0.5) + 权重衰减(10⁻⁵)",
    "梯度裁剪：max_norm = 5.0",
    "早停：patience = 7 epochs",
]
for i, c in enumerate(cfg):
    T(s, 0.8, 4.3 + i*0.27, 5.5, 0.25, c, 13, DARK)

# 右侧公式（全部绝对定位，无重叠）
B(s, 7.0, 1.35, 5.5, 0.35, "前向传播公式")
fy = 1.85

T(s, 7.0, fy, 5.5, 0.28, "单步 RNN 单元：", 13, DARK)
fy += 0.28
_, h = formula(s, 7.0, fy, "rnn1", r"h_t = \tanh\,(W_{hh}\!\cdot\!h_{t-1} + W_{xh}\!\cdot\!x_t + b_h)", 5.2, 13)
fy += h + 0.18

T(s, 7.0, fy, 5.5, 0.28, "双向拼接：", 13, DARK)
fy += 0.28
_, h = formula(s, 7.0, fy, "rnn2", r"h_t^{\mathrm{bi}} = \left[\,\overrightarrow{h_t}\;\|\;\overleftarrow{h_t}\,\right] \in \mathbb{R}^{512}", 5.2, 13)
fy += h + 0.18

T(s, 7.0, fy, 5.5, 0.28, "分类输出：", 13, DARK)
fy += 0.28
_, h = formula(s, 7.0, fy, "rnn3", r"\hat{y} = \mathrm{softmax}\,(W_{fc}\!\cdot\!h_{\mathrm{last}} + b_{fc})", 5.2, 13)
fy += h + 0.28

B(s, 7.0, fy, 5.5, 0.35, "固有局限：梯度消失")
T(s, 7.0, fy+0.42, 5.5, 0.25, "BPTT 中 ∂L/∂h₁ ∝ Π(W_hh)ᵏ → 指数衰减", 13, RED)
T(s, 7.0, fy+0.72, 5.5, 0.25, "→ 引出 LSTM 门控机制（组员 B 将介绍）", 13, DARK)

# 底部框
sh = s.slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.8))
sh.fill.solid(); sh.fill.fore_color.rgb = LIGHT; sh.line.color.rgb = ORANGE; sh.line.width = Pt(1)
T(s, 0.9, 6.4, 11.5, 0.3, "核心发现：5 倍数据量（16K → 80K），RNN 准确率从 81.79% 提升至 89.82%（+8.03 个百分点）", 15, BLUE, bold=True)
T(s, 0.9, 6.7, 11.5, 0.25, "数据规模的增长超过了任何架构改进带来的收益", 14, DARK)

# ═══════════════════════════════════════════════
# 第 7 页 · 结果与总结
# ═══════════════════════════════════════════════
s = Slide(6)
s.title("六、实验结果与讨论")

T(s, 3.0, 1.7, 7.5, 0.8, "RNN 测试准确率：81.79%  →  89.82%", 34, DARK, PP_ALIGN.CENTER, True)
T(s, 3.0, 2.6, 7.5, 0.5, "+8.03 个百分点（5 倍训练数据）", 24, GREEN, PP_ALIGN.CENTER, True)
L(s, 4.5, 3.3, 4.3)

B(s, 1.5, 3.7, 10, 0.5, "主要发现")
findings = [
    "8 个百分点的涨幅超过了 16K 小数据集上任何复杂模型的绝对性能。",
    "数据质量（预处理、否定词保护）直接影响所有下游模型的表现上限。",
    "RNN 达到 89.82% 验证了整个管线（分词、词向量、训练策略）的有效性。",
    "后续将由组员 B 介绍 BiLSTM 与 Attention 机制的架构设计。",
]
for i, f in enumerate(findings):
    T(s, 1.5, 4.25 + i*0.45, 10.3, 0.4, f"{i+1}.  {f}", 16, DARK)

# 交棒
sh = s.slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(6.5), Inches(6.3), Inches(0.6))
sh.fill.solid(); sh.fill.fore_color.rgb = BLUE; sh.line.fill.background()
tf = sh.text_frame
tf.paragraphs[0].text = "下一部分：BiLSTM 与 Attention 机制（组员 B）"
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
_set_font(tf.paragraphs[0], 17, True, WHITE)

out = os.path.join(BASE, "组员A_汇报PPT.pptx")
prs.save(out)
print(f"已生成：{out}")
