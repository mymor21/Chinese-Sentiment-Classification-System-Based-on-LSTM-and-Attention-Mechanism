"""
PPT 模板生成器 — 含 LaTeX 公式渲染支持
组员 B / C 直接替换 CONTENT 内容即可

依赖：pip install python-pptx matplotlib
"""
import os, io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# matplotlib 公式渲染
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ═══════════════════════════════════════════════
# 配色
# ═══════════════════════════════════════════════
DARK    = RGBColor(0x1E, 0x29, 0x3B)
PRIMARY = RGBColor(0x4F, 0x46, 0xE5)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xF8, 0xFA, 0xFC)
GRAY    = RGBColor(0x64, 0x74, 0x8B)
RED     = RGBColor(0xEF, 0x44, 0x44)
GREEN   = RGBColor(0x22, 0xC5, 0x5E)
ORANGE  = RGBColor(0xF5, 0x9E, 0x0B)
MUTED   = RGBColor(0x94, 0xA3, 0xB8)

# ═══════════════════════════════════════════════
# 公式渲染工具
# ═══════════════════════════════════════════════

# 预定义公式（项目中用到的所有公式）
FORMULAS = {}

def _render(formula_id, latex):
    """将 LaTeX 公式渲染为 PNG，存入 BytesIO 缓存"""
    fig, ax = plt.subplots(figsize=(0.01, 0.01))
    # 先用 mathtext 渲染看看实际尺寸
    # matplotlib 的 mathtext 支持常见 LaTeX：分数、上下标、希腊字母、求和等
    text = ax.text(0, 0, f"${latex}$", fontsize=16,
                   ha='left', va='bottom',
                   math_fontfamily='stix')
    fig.canvas.draw()
    bbox = text.get_window_extent(renderer=fig.canvas.get_renderer())
    w_inch = bbox.width / fig.dpi
    h_inch = bbox.height / fig.dpi
    plt.close(fig)

    # 以实际尺寸渲染
    fig, ax = plt.subplots(figsize=(max(w_inch + 0.15, 1), max(h_inch + 0.12, 0.3)))
    ax.axis('off')
    ax.text(0, 0.02, f"${latex}$", fontsize=16,
            ha='left', va='bottom', color='#1E293B',
            math_fontfamily='stix')
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=200, bbox_inches='tight',
                pad_inches=0.04, transparent=True)
    plt.close(fig)
    buf.seek(0)
    width_emu = int(w_inch * 914400)
    FORMULAS[formula_id] = (buf, width_emu)

def formula(slide, left, top, formula_id, latex):
    """将 LaTeX 公式渲染后插入到幻灯片指定位置"""
    if formula_id not in FORMULAS:
        _render(formula_id, latex)
    buf, w_emu = FORMULAS[formula_id]
    # 根据宽度等比例缩放
    pic = slide.shapes.add_picture(buf, Inches(left), Inches(top))
    scale = 914400 * 5.5 / w_emu  # 目标宽度 ~5.5 英寸
    pic.width = int(w_emu * min(scale, 1.0))
    pic.height = int(pic.height * min(scale, 1.0))
    return pic

# ═══════════════════════════════════════════════
# 工具函数
# ═══════════════════════════════════════════════

def add_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def add_title_bar(slide, title_text):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), prs.slide_width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.8)
    tf.margin_top = Inches(0.2)

def add_text(slide, left, top, width, height, text,
             size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Microsoft YaHei"
    p.alignment = align
    return tf

def add_para(tf, text, size=16, bold=False, color=DARK, before=4, after=2):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Microsoft YaHei"
    p.space_before = Pt(before)
    p.space_after = Pt(after)
    return p

def add_card(slide, x, y, w, h, border_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT
    shape.line.color.rgb = border_color
    shape.line.width = Pt(2)
    return shape

def add_divider(slide, y, w=3.3):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches((13.333 - w) / 2), Inches(y), Inches(w), Inches(0.03))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()

def add_page_num(slide, num):
    add_text(slide, 12.2, 7.0, 0.8, 0.4, str(num), size=12, color=GRAY, align=PP_ALIGN.RIGHT)

def make_table(slide, x, y, headers, rows, col_widths):
    for j, (cw, hdr) in enumerate(zip(col_widths, headers)):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(x + sum(col_widths[:j])), Inches(y), Inches(cw), Inches(0.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = PRIMARY
        shape.line.fill.background()
        tf2 = shape.text_frame
        tf2.paragraphs[0].text = hdr
        tf2.paragraphs[0].font.size = Pt(14)
        tf2.paragraphs[0].font.bold = True
        tf2.paragraphs[0].font.color.rgb = WHITE
        tf2.paragraphs[0].font.name = "Microsoft YaHei"
        tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
    for i, row in enumerate(rows):
        for j, (cw, val) in enumerate(zip(col_widths, row)):
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                Inches(x + sum(col_widths[:j])), Inches(y + 0.5 + i * 0.45),
                Inches(cw), Inches(0.45))
            shape.fill.solid()
            shape.fill.fore_color.rgb = LIGHT if i % 2 == 0 else WHITE
            shape.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
            shape.line.width = Pt(0.5)
            tf2 = shape.text_frame
            tf2.paragraphs[0].text = val
            tf2.paragraphs[0].font.size = Pt(14)
            tf2.paragraphs[0].font.color.rgb = DARK
            tf2.paragraphs[0].font.name = "Microsoft YaHei"
            tf2.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_compare_boxes(slide, y, left_title, left_lines, right_title, right_lines):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(y), Inches(5.6), Inches(max(len(left_lines), len(right_lines)) * 0.4 + 0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFE, 0xF2, 0xF2)
    shape.line.color.rgb = RED
    tf = add_text(slide, 1.1, y + 0.1, 5.0, 0.4, left_title, size=16, bold=True, color=RED)
    for t, c in left_lines:
        add_para(tf, t, size=15, color=c, before=4)

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(y), Inches(5.7), Inches(max(len(left_lines), len(right_lines)) * 0.4 + 0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF0, 0xFD, 0xF4)
    shape.line.color.rgb = GREEN
    tf = add_text(slide, 7.1, y + 0.1, 5.2, 0.4, right_title, size=16, bold=True, color=GREEN)
    for t, c in right_lines:
        add_para(tf, t, size=15, color=c, before=4)


# ═══════════════════════════════════════════════
# 第 1 页：封面
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_text(slide, 1, 2.5, 11.3, 1.2, "你的标题",
         size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, 1, 3.8, 11.3, 0.6, "副标题",
         size=22, color=MUTED, align=PP_ALIGN.CENTER)
add_divider(slide, 4.6)
add_text(slide, 1, 5.0, 11.3, 0.5, "组员 X：负责内容",
         size=18, color=WHITE, align=PP_ALIGN.CENTER)
add_page_num(slide, 1)

# ═══════════════════════════════════════════════
# 第 2 页：卡片布局示例
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "模型架构对比")

items = [
    ("BiLSTM", "双向 LSTM\n门控机制\n11.5M 参数", BLUE),
    ("Attention-LSTM", "加性 Self-Attention\n可解释性强\n11.8M 参数", PURPLE),
    ("CNN-BiLSTM", "多尺度卷积\n自动学 n-gram\n12.1M 参数", GREEN),
]
for i, (t, d, c) in enumerate(items):
    x = 0.8 + i * 4.2
    add_card(slide, x, 1.8, 3.8, 2.2, c)
    add_text(slide, x+0.3, 2.0, 3.2, 0.5, t, size=20, bold=True, color=c)
    for j, line in enumerate(d.split("\n")):
        add_text(slide, x+0.3, 2.6 + j*0.35, 3.2, 0.35, line, size=15, color=DARK)

add_page_num(slide, 2)

# ═══════════════════════════════════════════════
# 第 3 页：公式页（核心）
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "核心公式 — Attention 机制推导")

# 左侧架构
tf = add_text(slide, 0.8, 1.5, 5.5, 0.5, "架构总览", size=22, bold=True, color=PRIMARY)
add_para(tf, "", size=6)
add_para(tf, "Embedding(300d) → BiLSTM(2层)", size=15, color=DARK)
add_para(tf, "       ↓", size=15, color=GRAY, bold=True)
add_para(tf, "Self-Attention → 加权上下文向量", size=15, color=PRIMARY, bold=True)
add_para(tf, "       ↓", size=15, color=GRAY, bold=True)
add_para(tf, "Dropout → Linear → Softmax", size=15, color=DARK)
add_para(tf, "       ↓", size=15, color=GRAY, bold=True)
add_para(tf, "输出：正面 / 负面", size=15, color=GREEN, bold=True)
add_para(tf, "", size=10)
add_para(tf, "参数量：11.8M", size=16, bold=True, color=PRIMARY)

# 右侧公式 — 用 LaTeX 渲染
tf = add_text(slide, 7.0, 1.5, 5.5, 0.5, "核心公式", size=22, bold=True, color=PRIMARY)

add_para(tf, "加性 Attention（Bahdanau-style）：", size=15, bold=True, color=DARK, before=12)
formula(slide, 7.0, 2.2, "attn1", r"u_t = \tanh(W \cdot h_t)")

formula(slide, 7.0, 2.8, "attn2", r"\alpha_t = \mathrm{softmax}(u_t^{\top} \cdot u_w)")

formula(slide, 7.0, 3.4, "attn3", r"c = \sum_{t=1}^{T} \alpha_t \cdot h_t")

add_para(tf, "", size=6)
add_para(tf, "LSTM 门控公式：", size=15, bold=True, color=DARK, before=16)

formula(slide, 7.0, 4.3, "lstm1", r"f_t = \sigma(W_f \cdot [h_{t-1}, x_t] + b_f)")

formula(slide, 7.0, 4.85, "lstm2", r"i_t = \sigma(W_i \cdot [h_{t-1}, x_t] + b_i)")

# 底部优化策略
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.2))
shape.fill.solid()
shape.fill.fore_color.rgb = LIGHT
shape.line.color.rgb = ORANGE
tf = add_text(slide, 1.1, 5.9, 11.0, 0.4,
              "统一训练策略", size=16, bold=True, color=ORANGE)
add_para(tf, "Adam(lr=1e-3) + ReduceLROnPlateau(patience=3) + Dropout(0.5) + WeightDecay(1e-5) + GradClip(5.0) + EarlyStop(patience=7)",
         size=13, color=DARK)

add_page_num(slide, 3)

# ═══════════════════════════════════════════════
# 第 4 页：表格 + 结论
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "实验结果")

# 左：表格
add_text(slide, 0.8, 1.6, 5, 0.5, "测试集准确率", size=22, bold=True, color=PRIMARY)
make_table(slide, 0.8, 2.3,
    headers=["模型", "Acc", "F1"],
    rows=[
        ("RNN", "89.82%", "0.898"),
        ("BiLSTM", "92.44%", "0.924"),
        ("Attn-LSTM", "92.37%", "0.924"),
        ("CNN-BiLSTM", "92.12%", "0.921"),
        ("BERT", "94.66%", "0.947"),
    ],
    col_widths=[1.8, 1.2, 1.2])

# 右：核心发现
add_text(slide, 6.8, 1.6, 5.5, 0.5, "核心发现", size=22, bold=True, color=PRIMARY)
findings = [
    "1. 数据规模 > 模型复杂度",
    "2. BiLSTM 性价比最高（92.44%）",
    "3. Attention 短文本增益有限",
    "4. BERT 靠预训练知识拉开差距",
]
for i, f in enumerate(findings):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(2.3 + i * 0.7), Inches(5.5), Inches(0.55))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT
    shape.line.color.rgb = PRIMARY if i == 0 else RGBColor(0xE2,0xE8,0xF0)
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = f
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = (i == 0)
    tf.paragraphs[0].font.color.rgb = PRIMARY if i == 0 else DARK
    tf.paragraphs[0].font.name = "Microsoft YaHei"
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.08)

add_page_num(slide, 4)

# ═══════════════════════════════════════════════
# 第 5 页：总结 + 交棒
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "总结")

add_text(slide, 1, 2.0, 11.3, 1.0, "BiLSTM: 92.44%  vs  Attention-LSTM: 92.37%",
         size=32, bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_text(slide, 1, 3.0, 11.3, 0.6, "BiLSTM 是从零训练中性价比最高的模型",
         size=24, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_divider(slide, 3.8)

conclusions = [
    "1. LSTM 的门控机制有效缓解了 RNN 的梯度消失问题",
    "2. 双向建模让模型同时捕捉上下文信息",
    "3. Attention 提供可解释性，但在短文本上增益有限",
]
for i, c in enumerate(conclusions):
    add_text(slide, 1.5, 4.2 + i * 0.55, 10.3, 0.5, c, size=19, color=DARK)

# 交棒
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(3.5), Inches(6.5), Inches(6.3), Inches(0.6))
shape.fill.solid()
shape.fill.fore_color.rgb = PRIMARY
shape.line.fill.background()
tf = shape.text_frame
tf.paragraphs[0].text = "下面由组员 C 介绍实验结果与 Demo"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].font.name = "Microsoft YaHei"
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
add_page_num(slide, 5)

# ── 保存 ──────────────────────────────────────
prs.save("PPT模板_含公式.pptx")
print("已生成 PPT模板_含公式.pptx")
