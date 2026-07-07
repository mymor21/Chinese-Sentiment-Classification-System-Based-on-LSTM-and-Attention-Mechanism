"""
完整汇报PPT — 学术风格，优化布局，16页
"""
import os, io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import matplotlib; matplotlib.use("Agg")
import matplotlib.pyplot as plt

BASE = os.path.dirname(os.path.abspath(__file__))
R = os.path.join(BASE, "results")

# ── 配色 ──
NV = RGBColor(0x0F,0x17,0x2A); BL = RGBColor(0x25,0x62,0xEB)
LG = RGBColor(0xF1,0xF5,0xF9); WH = RGBColor(0xFF,0xFF,0xFF)
DK = RGBColor(0x1E,0x29,0x3B); GR = RGBColor(0x64,0x74,0x8B)
MU = RGBColor(0x94,0xA3,0xB8); RD = RGBColor(0xDC,0x26,0x26)
GN = RGBColor(0x16,0xA3,0x4A); OR = RGBColor(0xEA,0x58,0x0C)
BR = RGBColor(0xFE,0xF2,0xF2); BG = RGBColor(0xF0,0xFD,0xF4)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

# ── 公式 ──
FM = {}
def _RD(fid, latex, fs=15):
    fig, ax = plt.subplots(figsize=(0.01,0.01))
    t = ax.text(0,0,f"${latex}$",fontsize=fs,ha='left',va='bottom',math_fontfamily='stix')
    fig.canvas.draw()
    bb = t.get_window_extent(renderer=fig.canvas.get_renderer())
    w, h = bb.width/fig.dpi+0.10, bb.height/fig.dpi+0.08
    plt.close(fig)
    fig, ax = plt.subplots(figsize=(max(w,1.5),max(h,0.28)))
    ax.axis('off')
    ax.text(0,0.02,f"${latex}$",fontsize=fs,ha='left',va='bottom',color='#1E293B',math_fontfamily='stix')
    buf = io.BytesIO()
    fig.savefig(buf,format='png',dpi=200,bbox_inches='tight',pad_inches=0.03,transparent=True)
    plt.close(fig); buf.seek(0)
    FM[fid] = (buf,w,h)

def EQ(slide, left, top, fid, latex, max_w=5.5, fs=15):
    if fid not in FM: _RD(fid, latex, fs)
    buf, w, h = FM[fid]
    sld = slide.slide if hasattr(slide,'slide') else slide
    pic = sld.shapes.add_picture(buf, Inches(left), Inches(top))
    sc = min(max_w/w, 1.0)
    pic.width = int(w*914400*sc); pic.height = int(pic.height*sc)
    return pic, h*sc

# ── 页面 ──
class S:
    def __init__(self, num=None):
        self.slide = prs.slides.add_slide(prs.slide_layouts[6]); self.num = num
        bg = self.slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = WH
        bar = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.95))
        bar.fill.solid(); bar.fill.fore_color.rgb = NV; bar.line.fill.background()
        ln = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.95), prs.slide_width, Inches(0.02))
        ln.fill.solid(); ln.fill.fore_color.rgb = BL; ln.line.fill.background()
        if num: TX(self,12.3,7.05,0.7,0.3,str(num),10,GR,PP_ALIGN.RIGHT)
    def ti(self, text): TX(self,0.7,0.18,12,0.6,text,25,WH,PP_ALIGN.LEFT,True)

def _FT(p, sz=None, b=False, c=None):
    rp = p.font._rPr if hasattr(p,'font') else p._pPr
    for t, fn in [(qn('a:latin'),'Times New Roman'), (qn('a:ea'),'SimSun')]:
        e = rp.find(t)
        if e is None: e = rp.makeelement(t,{}); rp.insert(0,e)
        e.set('typeface',fn)
    if sz is not None: p.font.size = Pt(sz)
    if b: p.font.bold = True
    if c is not None: p.font.color.rgb = c

def TX(s, l, t, w, h, text, sz=14, c=DK, a=PP_ALIGN.LEFT, b=False):
    tb = s.slide.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = a; _FT(p,sz,b,c)
    return tf

def HD(s, l, t, w, h, text, c=BL, sz=17):
    TX(s,l,t,w,h,text,sz,c,PP_ALIGN.LEFT,True)

def LN(s, x, y, w):
    sh = s.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.015))
    sh.fill.solid(); sh.fill.fore_color.rgb = BL; sh.line.fill.background()

def CD(s, x, y, w, h, bdr):
    sh = s.slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = LG; sh.line.color.rgb = bdr; sh.line.width = Pt(1.5)

def IM(s, path, l, t, w, h=None):
    p = os.path.join(BASE, path)
    if os.path.exists(p):
        if h: s.slide.shapes.add_picture(p, Inches(l), Inches(t), Inches(w), Inches(h))
        else: s.slide.shapes.add_picture(p, Inches(l), Inches(t), Inches(w))

def HO(s, text):
    sh = s.slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5),Inches(6.55),Inches(6.3),Inches(0.5))
    sh.fill.solid(); sh.fill.fore_color.rgb = BL; sh.line.fill.background()
    tf = sh.text_frame; tf.paragraphs[0].text = text; _FT(tf.paragraphs[0],15,True,WH); tf.paragraphs[0].alignment = PP_ALIGN.CENTER


# ═══════════════════════════════════════════════
# 封面
# ═══════════════════════════════════════════════
s = S()
s.slide.background.fill.solid(); s.slide.background.fill.fore_color.rgb = NV
TX(s,1.2,1.6,11,1.2,"基于 LSTM 与 Attention 的中文情感分类系统",36,WH,PP_ALIGN.CENTER,True)
LN(s,5.2,3.1,3.0)
TX(s,1.2,3.4,11,0.4,"Chinese Sentiment Classification Based on LSTM and Attention Mechanism",15,MU,PP_ALIGN.CENTER)
TX(s,1.2,4.2,11,0.3,"自然语言处理  课程项目汇报",14,RGBColor(0xCB,0xD5,0xE1),PP_ALIGN.CENTER)
TX(s,1.2,5.0,11,0.3,"组员 A：数据管线 + RNN 基线    组员 B：模型架构设计    组员 C：实验结果 + Demo",13,MU,PP_ALIGN.CENTER)
TX(s,1.2,6.3,11,0.3,"2026 年 7 月",12,GR,PP_ALIGN.CENTER)

# ══════════════ A1: 问题定义 ══════════════
s = S(1); s.ti("一、研究问题与目标  （组员 A）")
HD(s,0.8,1.25,5,0.28,"任务定义")
TX(s,0.8,1.58,11,0.25,"给定中文评论文本，判断情感极性：正面 / 负面（二分类，y \u2208 {0, 1}）",14,DK)
HD(s,0.8,2.0,5,0.28,"研究目标")
TX(s,0.8,2.33,11,0.25,"对比五种模型架构（RNN \u2192 BiLSTM \u2192 Attention-LSTM \u2192 CNN-BiLSTM \u2192 BERT）的性能差异",14,DK)
TX(s,0.8,2.62,11,0.25,"探究数据规模（16K vs 80K）与模型复杂度对分类准确率的相对影响",14,DK)
HD(s,0.8,3.1,5,0.28,"损失函数")
_,h = EQ(s,0.8,3.5,"loss",r"\mathcal{L} = -\frac{1}{N}\sum_{i=1}^{N}\left[y_i\log\hat{p}_i + (1-y_i)\log(1-\hat{p}_i)\right]",11.5,14)
HD(s,0.8,4.35,5,0.28,"数据规模")
TX(s,0.8,4.68,11,0.25,"初版 16,000 条（酒店 + 外卖）\u2192 终版 80,000 条（+ 电商 10 类），正负各 40,000 条",14,DK)
HD(s,0.8,5.15,5,0.28,"对比模型")
TX(s,0.8,5.5,11,0.25,"RNN \u2192 BiLSTM \u2192 Attention-LSTM \u2192 CNN-BiLSTM \u2192 BERT（预训练基线，102M 参数）",14,BL,PP_ALIGN.LEFT,True)
HO(s,"下一部分：数据获取与预处理")

# ══════════════ A2: 数据获取 ══════════════
s = S(2); s.ti("二、数据获取与预处理  （组员 A）")
ds = [("ChnSentiCorp","9,600 条","酒店评论","HuggingFace",BL),
      ("waimai_10k","11,987 条","外卖评论","阿里云 OSS",OR),
      ("online_shopping","62,774 条","电商 10 类别","本地 CSV",GN)]
for i,(n,sz,dom,src,cl) in enumerate(ds):
    x = 0.85 + i*4.15; CD(s,x,1.3,3.7,1.5,cl)
    TX(s,x+0.25,1.38,3.0,0.3,n,16,cl,PP_ALIGN.LEFT,True)
    TX(s,x+0.25,1.72,3.0,0.24,sz,14,DK); TX(s,x+0.25,2.0,3.0,0.22,dom,12,GR)
    TX(s,x+0.25,2.26,3.0,0.22,f"来源：{src}",11,MU)
TX(s,1.5,3.15,10.3,0.4,"合并 84,207 条  \u2192  平衡采样  \u2192  80,000 条（每类各 40,000 条）",19,BL,PP_ALIGN.CENTER,True)
IM(s,"results/analysis/sentiment_distribution.png",3.5,3.7,6.2,2.6)
TX(s,0.8,6.55,11.5,0.22,"预处理：清洗 \u2192 jieba 分词 \u2192 停用词过滤（保留否定词）\u2192 词表 29,343 \u2192 Word2Vec(300d) \u2192 前向 Padding(67)",12,DK)

# ══════════════ A3: 预处理改进 ══════════════
s = S(3); s.ti("三、预处理关键改进  （组员 A）")
steps = ["原始文本","清洗","分词\n(jieba)","停用词\n过滤","词表\n(29K)","Word2Vec\n(300d)","Padding\n(67)"]
for i,st in enumerate(steps):
    x = 0.65 + i*1.78
    sh = s.slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.3), Inches(1.45), Inches(0.62))
    sh.fill.solid(); sh.fill.fore_color.rgb = BL if i < 6 else NV; sh.line.fill.background()
    tf = sh.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.text = st; _FT(p,10,True,WH); p.alignment = PP_ALIGN.CENTER
    if i < 6: TX(s,x+1.45,1.35,0.28,0.25,"\u25b6",11,GR,PP_ALIGN.CENTER)

# 对比
CD(s,0.6,2.25,5.8,2.5,RD); s.slide.shapes[-1].fill.fore_color.rgb = BR
TX(s,0.9,2.35,5.2,0.28,"原方案：启发式单字过滤（len(w) < 2）",14,RD,PP_ALIGN.LEFT,True)
TX(s,0.9,2.72,5.2,0.28,"所有单字词直接丢弃：好、坏、差、不、低、贵等情感关键词被误删",13,DK)
TX(s,0.9,3.08,5.2,0.28,"\u300e服务不热情\u300f \u2192 \u300e服务热情\u300f（语义反转！）",13,RD)
TX(s,0.9,3.42,5.2,0.28,"\u300e性价比很低\u300f \u2192 \u300e性价比\u300f（否定信息丢失）",13,RD)
TX(s,0.9,3.8,5.2,0.28,"\u300e一问三不知\u300f \u2192 \u2192（全部删除）",13,RD)

CD(s,6.8,2.25,5.9,2.5,GN); s.slide.shapes[-1].fill.fore_color.rgb = BG
TX(s,7.1,2.35,5.2,0.28,"改进方案：精确停用词表过滤",14,GN,PP_ALIGN.LEFT,True)
_,h = EQ(s,7.1,2.8,"sw",r"W' = \{\,w \in W \mid w \notin S\,\}",4.5,13)
TX(s,7.1,3.42,5.2,0.28,"删除：虚词、语气词、标点、数字、英文",13,DK)
TX(s,7.1,3.78,5.2,0.28,"保留：否定词(不/没/非) + 程度词(很/太/最) + 单字实词",13,GN)
TX(s,7.1,4.15,5.2,0.28,"效果：所有模型准确率提升 1~4 个百分点",14,GN,PP_ALIGN.LEFT,True)

TX(s,0.8,5.05,11.5,0.22,"参数：词表 29,343 词  |  Skip-gram, 300d, window=5, 30 epochs  |  max_seq_len=67 (95%分位数), 前向 Padding",12,DK)
TX(s,0.8,5.32,11.5,0.22,"覆盖率 100%  |  Train/Val/Test = 70/15/15（分层采样）",12,DK)
HO(s,"下一部分：Word2Vec 词向量训练")

# ══════════════ A4: Word2Vec ══════════════
s = S(4); s.ti("四、词向量训练：Skip-gram 与负采样  （组员 A）")
HD(s,0.8,1.25,5.5,0.28,"训练配置")
for i,(k,v) in enumerate([("架构","Skip-gram (sg=1)"),("维度","300"),("上下文窗口","5"),("训练轮数","30"),("最低词频","min_count=3"),("负采样","K=5")]):
    TX(s,0.8,1.62+i*0.28,5.5,0.24,f"  {k}：{v}",13,DK)
HD(s,0.8,3.5,5.5,0.28,"语义空间")
TX(s,0.8,3.82,5.5,0.24,"语义相近的词在向量空间中距离更近",13,DK)
TX(s,0.8,4.1,5.5,0.24,"余弦相似度度量语义相关性",13,DK)

HD(s,7.0,1.25,5.5,0.28,"目标函数")
y = 1.72; TX(s,7.0,y,5.5,0.22,"Skip-gram：给定中心词 c，预测上下文词 o",12,DK); y += 0.26
_,h = EQ(s,7.0,y,"w2v",r"\max\ \sum_{t=1}^{T}\ \sum_{-c\leq j\leq c,\,j\neq 0}\ \log P(w_{t+j}\mid w_t)",5.5,12); y += h+0.15
TX(s,7.0,y,5.5,0.22,"负采样近似（将 |V| 分类降为 K+1 个二分类）：",12,DK); y += 0.26
_,h = EQ(s,7.0,y,"ns",r"\max\ \log\sigma(u_o^{\top}v_c) + \sum_{k=1}^{K}\mathbb{E}_{w_k\sim P_n}\left[\log\sigma(-u_k^{\top}v_c)\right]",5.5,12)
y += h+0.18; HD(s,7.0,y,5.5,0.28,"嵌入矩阵")
TX(s,7.0,y+0.32,5.5,0.22,"E \u2208 R^{V \u00d7 300}，V=29,343，预训练权重作为 Embedding 初始化",12,DK)

IM(s,"results/analysis/wordclouds.png",0.8,5.15,11.7,2.1)
HO(s,"下一部分：RNN 基线模型")

# ══════════════ A5: RNN ══════════════
s = S(5); s.ti("五、RNN 基线模型  （组员 A）")
HD(s,0.8,1.25,5.5,0.28,"模型架构")
TX(s,0.8,1.58,5.5,0.24,"Input(67) \u2192 Embedding(300d, 预训练) \u2192 2-layer BiRNN(h=256)",13,DK)
TX(s,0.8,1.86,5.5,0.24,"Last Hidden Concat \u2192 512d \u2192 Dropout(0.5) \u2192 Linear(2) \u2192 Softmax",13,DK)
TX(s,0.8,2.22,5.5,0.26,"参数量：9,483,862（\u2248 950 万，五个模型中最轻量）",14,BL,PP_ALIGN.LEFT,True)
HD(s,0.8,2.7,5.5,0.28,"训练配置")
TX(s,0.8,3.02,5.5,0.24,"Adam(lr=1e-3) + ReduceLROnPlateau(patience=3, factor=0.5)",13,DK)
TX(s,0.8,3.3,5.5,0.24,"Dropout(0.5) + WeightDecay(1e-5) + GradClip(5.0) + EarlyStop(7)",13,DK)

HD(s,7.0,1.25,5.5,0.28,"核心公式")
fy = 1.72; TX(s,7.0,fy,5.5,0.22,"RNN 单步前向传播：",12,DK); fy += 0.26
_,h = EQ(s,7.0,fy,"rnn1",r"h_t = \tanh\,(W_{hh} \cdot h_{t-1} + W_{xh} \cdot x_t + b_h)",5.2,12); fy += h+0.16
TX(s,7.0,fy,5.5,0.22,"双向拼接：",12,DK); fy += 0.26
_,h = EQ(s,7.0,fy,"rnn2",r"h_t^{\mathrm{bi}} = [\,\overrightarrow{h_t}\;\|\;\overleftarrow{h_t}\,] \in \mathbb{R}^{512}",5.2,12); fy += h+0.16
TX(s,7.0,fy,5.5,0.22,"分类输出：",12,DK); fy += 0.26
_,h = EQ(s,7.0,fy,"rnn3",r"\hat{y} = \mathrm{softmax}\,(W_{fc} \cdot h_{\mathrm{last}} + b_{fc})",5.2,12); fy += h+0.25
HD(s,7.0,fy,5.5,0.28,"局限：梯度消失")
TX(s,7.0,fy+0.32,5.5,0.22,"BPTT: \u2202L/\u2202h\u2081 \u221d \u03a0 (W_hh)^k \u2192 指数衰减，须引入 LSTM 门控",12,RD)

CD(s,0.6,6.2,12.1,0.65,OR); TX(s,0.85,6.28,11.5,0.26,"核心数据：5\u00d7 数据（16K \u2192 80K），RNN 准确率 81.79% \u2192 89.82%（+8.03 pp）\u2014\u2014 数据规模的增长超过任何架构改进的收益",14,BL,PP_ALIGN.LEFT,True)
HO(s,"下一部分：LSTM 门控机制  （组员 B）")

# ══════════════ B1: LSTM ══════════════
s = S(6); s.ti("六、LSTM 门控机制  （组员 B）")
HD(s,0.8,1.25,5.5,0.28,"从 RNN 到 LSTM")
TX(s,0.8,1.6,5.5,0.5,"RNN 的 BPTT 训练中梯度需连乘转移矩阵 W_hh 的幂。若最大特征值 < 1，梯度指数衰减，早期信息完全丢失。这是结构性缺陷，无法通过调参解决。",13,DK)
TX(s,0.8,2.35,5.5,0.24,"LSTM 引入细胞状态 C_t 和三个门控单元，用加法路径替代连乘：",13,DK)
TX(s,0.8,2.65,5.5,0.24,"遗忘门 f_t：\u03c3(W_f\u00b7[h_{t-1},x_t]+b_f)  \u2014 决定丢弃多少旧记忆",13,DK)
TX(s,0.8,2.95,5.5,0.24,"输入门 i_t：\u03c3(W_i\u00b7[h_{t-1},x_t]+b_i)  \u2014 决定写入多少新信息",13,DK)
TX(s,0.8,3.25,5.5,0.24,"输出门 o_t：\u03c3(W_o\u00b7[h_{t-1},x_t]+b_o)  \u2014 决定输出多少信息",13,DK)
TX(s,0.8,3.7,5.5,0.28,"关键：当 f_t \u2248 1 时，C_t \u2248 C_{t-1} + i_t\u2299tanh(\u00b7)，梯度沿此加法路径几乎无损传递",13,DK)
TX(s,0.8,4.15,5.5,0.26,"模型参数：BiLSTM, 2 层, 每方向 256d \u2192 拼接 512d, 11.5M 参数",14,BL,PP_ALIGN.LEFT,True)

HD(s,7.0,1.25,5.5,0.28,"LSTM 核心公式")
fy = 1.72
for fid, latex in [
    ("lstm1", r"f_t = \sigma(W_f \cdot [h_{t-1}, x_t] + b_f)"),
    ("lstm2", r"i_t = \sigma(W_i \cdot [h_{t-1}, x_t] + b_i)"),
    ("lstm3", r"\tilde{C}_t = \tanh(W_C \cdot [h_{t-1}, x_t] + b_C)"),
    ("lstm4", r"C_t = f_t \odot C_{t-1} + i_t \odot \tilde{C}_t"),
    ("lstm5", r"o_t = \sigma(W_o \cdot [h_{t-1}, x_t] + b_o),\quad h_t = o_t \odot \tanh(C_t)"),
]:
    _,h = EQ(s,7.0,fy,fid,latex,5.2,13); fy += h+0.13
TX(s,7.0,fy,5.5,0.22,"\u03c3 \u2208 (0,1) 作为门控系数，\u2299 逐元素乘法，[h_{t-1}, x_t] 拼接操作",11,GR)
HO(s,"下一部分：Attention 机制")

# ══════════════ B2: Attention ══════════════
s = S(7); s.ti("七、加性 Self-Attention 机制  （组员 B）")
HD(s,0.8,1.25,5.5,0.28,"动机")
TX(s,0.8,1.6,5.5,0.45,"BiLSTM 输出所有时间步的隐状态，但每个词对情感判断的贡献并不相同。Attention 让模型自动学习哪些词更关键，并为决策提供可解释性。",13,DK)
HD(s,0.8,2.3,5.5,0.28,"计算流程（Bahdanau 风格）")
for i,st in enumerate([
    "Step 1: u_t = tanh(W\u00b7h_t)  \u2014 对隐状态做非线性投影",
    "Step 2: score_t = u_t^T\u00b7u_w  \u2014 与可学习向量计算相似度",
    "Step 3: \u03b1_t = softmax(score_t)  \u2014 归一化为概率分布",
    "Step 4: c = \u03a3 \u03b1_t\u00b7h_t  \u2014 加权求和得上下文向量",
]):
    TX(s,0.8,2.65+i*0.3,5.5,0.26,st,13,DK)
TX(s,0.8,4.05,5.5,0.22,"Mask: padding 位置 score = -\u221e，\u03b1 = 0",11,GR)
TX(s,0.8,4.5,5.5,0.22,"W \u2208 R^{512\u00d7512}, u_w \u2208 R^{512} 均可学习",11,GR)

HD(s,7.0,1.25,5.5,0.28,"数学表达")
fy = 1.72
_,h = EQ(s,7.0,fy,"attn1",r"u_t = \tanh(W \cdot h_t)",5.2,15); fy += h+0.16
_,h = EQ(s,7.0,fy,"attn2",r"\alpha_t = \frac{\exp(u_t^{\top} u_w)}{\sum_j \exp(u_j^{\top} u_w)}",5.2,15); fy += h+0.16
_,h = EQ(s,7.0,fy,"attn3",r"c = \sum_{t=1}^{T} \alpha_t \cdot h_t",5.2,15)
fy += h+0.25
TX(s,7.0,fy,5.5,0.22,"\u03b1_t 可可视化 \u2192 注意力热力图，模型决策可解释",13,GN,PP_ALIGN.LEFT,True)
TX(s,7.0,fy+0.28,5.5,0.22,"参数量：Attention-LSTM 11.8M",13,DK)

# ══════════════ B3: CNN-BiLSTM ══════════════
s = S(8); s.ti("八、CNN-BiLSTM 多尺度特征  （组员 B）")
HD(s,0.8,1.25,11,0.28,"设计思路")
TX(s,0.8,1.58,6.5,0.38,"在 Embedding 和 BiLSTM 之间插入三路并行 1D 卷积。卷积核通过训练自动发现改变情感极性的局部 n-gram 模式，无需预设否定词列表。",13,DK)
HD(s,0.8,2.2,5.5,0.28,"三路卷积核")
TX(s,0.8,2.52,5.5,0.24,"kernel = 2 \u2192 二元组：不热情、很低、好吃",13,DK)
TX(s,0.8,2.8,5.5,0.24,"kernel = 3 \u2192 三元组：性价比很低、很不错",13,DK)
TX(s,0.8,3.08,5.5,0.24,"kernel = 4 \u2192 四元组：物超所值、一问三不知",13,DK)
TX(s,0.8,3.5,5.5,0.24,"每路 100 卷积核 \u2192 拼接后 Linear 投影 \u2192 输入 BiLSTM",13,DK)
HD(s,0.8,3.95,5.5,0.28,"优势与参数")
TX(s,0.8,4.27,5.5,0.24,"自动学习局部否定模式，无需人工维护否定词词典",13,DK)
TX(s,0.8,4.55,5.5,0.24,"参数量：12.1M（五个从零训练模型中最大）",14,BL,PP_ALIGN.LEFT,True)

HD(s,7.0,1.25,5.5,0.28,"架构流程")
fy = 1.72
for line in ["Embedding(300d, 预训练)","  \u2193","3 \u00d7 Conv1D (kernel=2,3,4)","  \u2193  Concat + Linear Project","BiLSTM(2-layer, h=256) + Attention","  \u2193","Dropout(0.5) \u2192 Linear(2) \u2192 Softmax"]:
    c = GR if "\u2193" in line else DK; sz = 12 if "\u2193" in line else 13
    TX(s,7.0,fy,5.5,0.25,line,sz,c); fy += 0.27
HO(s,"下一部分：统一训练策略")

# ══════════════ B4: 训练策略 ══════════════
s = S(9); s.ti("九、统一训练策略  （组员 B）")
HD(s,0.8,1.25,11,0.28,"从零训练四模型统一配置（RNN / LSTM / Attention-LSTM / CNN-BiLSTM）")
specs = [("损失函数","CrossEntropyLoss","标准二分类"),("优化器","Adam (lr = 10\u207b\u00b3)","自适应学习率"),("学习率调度","ReduceLROnPlateau (patience=3)","验证停滞 \u2192 lr \u00d7 0.5"),("正则化","Dropout(0.5) + WeightDecay(10\u207b\u2075)","双重防过拟合"),("梯度裁剪","max_norm = 5.0","防止梯度爆炸"),("早停","patience = 7 epochs","恢复历史最优模型"),("批次大小","64","GPU 显存与收敛平衡"),("最大 Epoch","50（通常早停触发）","理论上限")]
for i,(n,v,d) in enumerate(specs):
    y = 1.62 + i*0.34
    TX(s,0.8,y,2.8,0.3,f"  {n}",13,DK,PP_ALIGN.LEFT,True)
    TX(s,3.3,y,4.2,0.3,v,13,DK)
    TX(s,7.5,y,4.5,0.3,d,12,GR)
TX(s,0.8,4.55,11.5,0.24,"BERT 例外：AdamW (lr = 2\u00d710\u207b\u2075), batch_size = 16, max_epochs = 5, max_length = 128",12,GR)
TX(s,0.8,4.9,11.5,0.24,"四个从零训练模型：\u223c94% 参数在 Embedding 层（Word2Vec 预训练），RNN/LSTM 核心仅 \u223c0.6M 可训参数",12,DK)
HD(s,0.8,5.35,5,0.28,"训练流程")
TX(s,0.8,5.68,11.5,0.24,"for epoch in range(max_epochs):  train \u2192 grad_clip \u2192 evaluate(val) \u2192 scheduler.step \u2192 check early_stop",12,DK)
HO(s,"下一部分：实验结果  （组员 C）")

# ══════════════ C1: 全模型对比 ══════════════
s = S(10); s.ti("十、五模型实验结果对比  （组员 C）")
hdrs = ["模型","16K Acc","80K Acc","\u0394","参数量","特点"]
cw = [1.5,1.1,1.1,0.8,1.1,4.8]
rows = [("RNN","81.79%","89.82%","+8.03","9.5M","基线，最轻量"),
        ("BiLSTM","87.46%","92.44%","+4.98","11.5M","从零训练中性价比最高"),
        ("Attn-LSTM","87.71%","92.37%","+4.66","11.8M","注意力加权，可解释"),
        ("CNN-BiLSTM","87.83%","92.12%","+4.29","12.1M","自动学习局部 n-gram"),
        ("BERT","\u2014","94.66%","\u2014","102M","预训练知识天花板")]
tx, ty = 1.3, 1.5
for j,(w,hdr) in enumerate(zip(cw,hdrs)):
    TX(s,tx+sum(cw[:j]),ty,w,0.36,hdr,11,BL,PP_ALIGN.CENTER,True)
for i,row in enumerate(rows):
    ry = ty + 0.38 + i*0.36
    for j,(w,val) in enumerate(zip(cw,row)):
        c = GN if i == 0 else DK
        TX(s,tx+sum(cw[:j]),ry,w,0.34,val,11,c,PP_ALIGN.CENTER)
    if i < len(rows)-1:
        sh = s.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(tx), Inches(ry+0.35), Inches(sum(cw)), Inches(0.006))
        sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0xE2,0xE8,0xF0); sh.line.fill.background()

HD(s,1.3,3.75,10,0.3,"核心发现")
for i,f in enumerate([
    "1. 数据规模效应 > 架构改进：5\u00d7 数据令 RNN 涨幅 +8.03 pp，超过 16K 下任何复杂模型的绝对性能。",
    "2. BiLSTM 是从零训练的最优选择（92.44%, 11.5M），与 Attention-LSTM 差距仅 0.07 pp。",
    "3. BERT 94.66% 展示预训练范式的性能上限，2.2 pp 差距来自外部语料知识，非架构本身优势。"
]):
    TX(s,1.3,4.1+i*0.38,10.5,0.34,f,14,DK)
IM(s,"results/evaluation/model_comparison.png",1.0,5.35,5.5,1.85)
IM(s,"results/training/loss_curves.png",6.8,5.35,5.5,1.85)
HO(s,"下一部分：可视化分析")

# ══════════════ C2a: 混淆矩阵 + 注意力热力图 ══════════════
s = S(11); s.ti("十一、混淆矩阵与注意力可视化  （组员 C）")
IM(s,"results/evaluation/confusion_matrices.png",0.5,1.3,6.3,2.9)
TX(s,0.5,4.3,6.3,0.25,"图 1：五模型混淆矩阵（行 = 真实标签，列 = 预测标签）",11,GR,PP_ALIGN.CENTER)
IM(s,"results/analysis/attention_heatmaps.png",6.6,1.3,6.2,2.9)
TX(s,6.6,4.3,6.2,0.25,"图 2：Attention-LSTM 注意力权重热力图",11,GR,PP_ALIGN.CENTER)
HD(s,0.8,4.8,11,0.28,"分析要点")
TX(s,0.8,5.15,5.5,0.24,"混淆矩阵：所有模型正负类分类均衡，无明显偏向",13,DK)
TX(s,0.8,5.42,5.5,0.24,"注意力：模型能准确聚焦情感关键词（如「好」「差」）",13,DK)
TX(s,6.8,5.15,5.5,0.24,"BERT 混淆矩阵最干净，误分类案例极少",13,DK)
TX(s,6.8,5.42,5.5,0.24,"RNN 基线在长文本上的混淆多于短文本",13,DK)
HO(s,"下一部分：嵌入可视化与词云")

# ══════════════ C2b: t-SNE + 词云 ══════════════
s = S(12); s.ti("十二、嵌入可视化与情感词云  （组员 C）")
IM(s,"results/analysis/tsne_embeddings.png",0.5,1.3,6.3,2.9)
TX(s,0.5,4.3,6.3,0.25,"图 3：t-SNE 嵌入可视化（Attention-LSTM 特征空间）",11,GR,PP_ALIGN.CENTER)
IM(s,"results/analysis/wordclouds.png",6.6,1.3,6.2,2.9)
TX(s,6.6,4.3,6.2,0.25,"图 4：情感区分词云（左 = 负面特征词，右 = 正面特征词）",11,GR,PP_ALIGN.CENTER)
HD(s,0.8,4.8,11,0.28,"分析要点")
TX(s,0.8,5.15,5.5,0.24,"t-SNE：正面/负面评论在高维特征空间中已形成可分簇",13,DK)
TX(s,0.8,5.42,5.5,0.24,"簇边界清晰说明 Attention-LSTM 学到了有效情感表征",13,DK)
TX(s,6.8,5.15,5.5,0.24,"词云：跨类共性过滤排除了「酒店」「外卖」等领域通用词",13,DK)
TX(s,6.8,5.42,5.5,0.24,"「差」「好」「烂」「赞」等真正情感区分词得到突显",13,DK)
HO(s,"下一部分：交互式 Demo")

# ══════════════ C3: Demo ══════════════
s = S(13); s.ti("十三、交互式 Demo 与工程优化  （组员 C）")
items = [
    ("置信度三档分级","\u2265 85%：正常预测  |  65~85%：置信度偏低  |  < 65%：显示 [不确定]","按行业标准设定双阈值，避免低置信度错误引导"),
    ("BERT 后台异步加载","4 个轻量模型（RNN/LSTM/Attention-LSTM/CNN-BiLSTM）5 秒内就绪","BERT 约 50 秒后台加载，页面秒开不阻塞，前端轮询 /api/status"),
    ("词云跨类共性过滤","min/max > 0.25 \u2192 排除领域通用词  |  正/负比值 > 1.5 \u2192 真正情感区分词","词云不再被「酒店」「房间」等主题词占据，聚焦情感特征"),
    ("训练-推理一致性","app.py 与 preprocess.py 共用同一套 tokenize 函数","消除训练与部署环境的预处理差异，避免推理阶段性能下降"),
]
for i,(title,detail,note) in enumerate(items):
    y = 1.3 + i*1.38
    CD(s,0.6,y,12.1,1.22,BL)
    TX(s,0.85,y+0.06,11.5,0.26,title,16,BL,PP_ALIGN.LEFT,True)
    TX(s,0.85,y+0.38,11.5,0.24,detail,13,DK)
    TX(s,0.85,y+0.7,11.5,0.24,f"\u2192  {note}",12,GN)
TX(s,0.8,6.78,11.5,0.2,"技术栈：Flask + Chart.js + HTML/CSS  |  启动：python app.py \u2192 http://localhost:5000",11,GR)

# ══════════════ C4: 总结 ══════════════
s = S(14); s.ti("十四、总结与展望")
TX(s,2.0,1.5,9.5,0.6,"核心结论",27,BL,PP_ALIGN.CENTER,True)
LN(s,4.5,2.2,4.3)
conclusions = [
    "1. 数据规模 > 模型复杂度 \u2014\u2014 多喂高质量数据是最有效的提升手段（+8 pp > 任何架构改进）。",
    "2. 预处理质量决定模型上限 \u2014\u2014 停用词精确控制让所有模型统一提升 1~4 个百分点。",
    "3. BiLSTM 是从零训练的最优选择 \u2014\u2014 92.44% 准确率，参数量适中（11.5M），训练推理高效。",
    "4. 引入预训练是突破 92% 瓶颈的最直接路径 \u2014\u2014 BERT 94.66% 来自海量外部语料知识。",
]
for i,c in enumerate(conclusions):
    TX(s,1.5,2.6+i*0.52,10.3,0.46,c,16,DK)

CD(s,3.5,6.2,6.3,0.55,NV); s.slide.shapes[-1].fill.fore_color.rgb = NV
tf = s.slide.shapes[-1].text_frame; tf.paragraphs[0].text = "谢谢！欢迎提问"
_FT(tf.paragraphs[0],19,True,WH); tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# ── 保存 ──
out = os.path.join(BASE,"完整汇报PPT.pptx"); prs.save(out)
print(f"已生成：{out}")
